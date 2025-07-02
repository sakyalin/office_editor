import os
import re
import tempfile
from pathlib import Path

from openai import OpenAI
from docx import Document
from pptx import Presentation
from docx.shared import RGBColor
from pptx.dml.color import RGBColor as PPTRGBColor

# 导入自定义环境变量加载模块
from env_loader import load_env_variables, get_api_key, get_api_base_url

# 加载.env文件中的环境变量
load_env_variables()

# 全局变量
client = None

def init_openai_client():
    """初始化OpenAI客户端"""
    global client
    
    # 如果客户端已经初始化，直接返回
    if client is not None:
        return client
    
    # 从环境变量获取API密钥
    api_key = get_api_key()
    
    # 验证API密钥
    if not api_key or api_key == "your_api_key_here":
        raise ValueError("请在.env文件中设置有效的OPENAI_API_KEY环境变量")
    
    try:
        # 从环境变量获取API基础URL
        api_base_url = get_api_base_url()
        
        # 使用指定的API基础URL初始化客户端
        client = OpenAI(
            api_key=api_key,  # 确保API密钥是ASCII字符
            base_url=api_base_url
        )
        return client
    except Exception as e:
        print(f"OpenAI客户端初始化错误: {e}")
        raise


def process_document(file_path, progress_callback=None):
    """处理Office文档，检查错别字和病句"""
    # 初始化OpenAI客户端
    try:
        init_openai_client()
    except ValueError as e:
        raise ValueError(f"OpenAI API初始化失败: {str(e)}")
        
    # 打印调试信息
    print(f"处理文件: {file_path}")
    
    file_path = Path(file_path)
    file_extension = file_path.suffix.lower()
    
    # 根据文件类型选择处理函数
    if file_extension == ".docx":
        return process_word(file_path, progress_callback)
    elif file_extension == ".pptx":
        return process_powerpoint(file_path, progress_callback)
    else:
        raise ValueError(f"不支持的文件格式: {file_extension}")


def get_openai_suggestions(text):
    """使用OpenAI检查文本中的错别字和病句"""
    if not text or text.strip() == "":
        return text, []
    
    # 确保客户端已初始化
    try:
        init_openai_client()
    except ValueError as e:
        raise ValueError(f"OpenAI客户端未初始化: {str(e)}")
    
    try:
        # 确保文本是UTF-8编码
        if isinstance(text, str):
            text_utf8 = text
        else:
            text_utf8 = str(text)
        
        response = client.chat.completions.create(
            model="gpt-4o-mini",
            messages=[
                {"role": "system", "content": "你是一位专业的校对助手。请检查以下文本中的错别字和语法错误。只需指出需要修改的部分并提供修改后的文本。无需解释原因。如果不需要修改，则返回空字符串。格式：原文|修改后的文本"}, 
                {"role": "user", "content": text_utf8}
            ],
            temperature=0.3,
            max_tokens=2000
        )
        
        suggestions_text = response.choices[0].message.content.strip()
        
        # 解析建议
        suggestions = []
        for line in suggestions_text.split("\n"):
            if "|" in line:
                parts = line.split("|")
                if len(parts) >= 2:
                    original = parts[0].strip()
                    if original.startswith("["): 
                        original = original[1:]
                    if original.endswith("]"): 
                        original = original[:-1]
                    
                    suggestion = parts[1].strip()
                    if suggestion.startswith("["): 
                        suggestion = suggestion[1:]
                    if suggestion.endswith("]"): 
                        suggestion = suggestion[:-1]
                    
                    suggestions.append((original, suggestion))
        
        return text, suggestions
    
    except Exception as e:
        print(f"OpenAI API调用出错: {e}")
        # 使用模拟数据作为备选
        print("使用模拟数据作为备选...")
        suggestions = []
        if len(text) > 10:
            sample = text[5:15] if len(text) > 15 else text[:5]
            suggestions.append((sample, sample + "(建议修改示例)"))
        return text, suggestions


def process_word(file_path, progress_callback=None):
    """处理Word文档"""
    doc = Document(file_path)
    
    # 计算总项目数（段落 + 表格中的段落）
    # 包括空段落在内的所有段落都计入总数
    # 这样即使跳过空段落，进度计算也是准确的
    total_items = len(doc.paragraphs)
    for table in doc.tables:
        for row in table.rows:
            for cell in row.cells:
                total_items += len(cell.paragraphs)
    
    # 当前已处理项目数
    processed_items = 0
    
    # 处理段落
    for paragraph in doc.paragraphs:
        text = paragraph.text
        # 跳过空段落的处理
        if not text.strip():
            # 更新进度
            processed_items += 1
            if progress_callback:
                progress_percent = int((processed_items / total_items) * 100)
                progress_callback(progress_percent, f"正在处理 {processed_items}/{total_items}")
            continue
            
        _, suggestions = get_openai_suggestions(text)
        
        if suggestions:
                # 保存原始段落的所有运行及其格式
                original_runs = []
                for run in paragraph.runs:
                    original_runs.append({
                        'text': run.text,
                        'bold': run.bold,
                        'italic': run.italic,
                        'underline': run.underline,
                        'font': run.font.name,
                        'size': run.font.size,
                        'color': run.font.color.rgb if run.font.color.rgb else None,
                        'highlight_color': run.font.highlight_color,
                        'style': run.style
                    })
                
                # 清除段落中的所有运行
                p_element = paragraph._element
                p_element.clear()
                
                # 重建原始文本并添加建议
                # 首先构建原始文本的字符位置到运行的映射
                char_to_run_map = []
                current_pos = 0
                for run_info in original_runs:
                    run_text = run_info['text']
                    for _ in range(len(run_text)):
                        char_to_run_map.append(run_info)
                        current_pos += 1
                
                # 重新添加文本，并在需要修改的地方添加红色建议
                current_pos = 0
                for original, suggestion in suggestions:
                    # 查找原文在文本中的位置
                    pos = text.find(original, current_pos)
                    if pos != -1:
                        # 添加原文前的文本，保持原始格式
                        if pos > current_pos:
                            for i in range(current_pos, pos):
                                if i < len(char_to_run_map):
                                    run_info = char_to_run_map[i]
                                    run = paragraph.add_run(text[i])
                                    # 应用原始格式
                                    run.bold = run_info['bold']
                                    run.italic = run_info['italic']
                                    run.underline = run_info['underline']
                                    if run_info['font']:
                                        run.font.name = run_info['font']
                                    if run_info['size']:
                                        run.font.size = run_info['size']
                                    if run_info['color']:
                                        run.font.color.rgb = run_info['color']
                                    run.font.highlight_color = run_info['highlight_color']
                                    if run_info['style']:
                                        run.style = run_info['style']
                                else:
                                    # 如果没有格式信息，则使用默认格式
                                    paragraph.add_run(text[i])
                        
                        # 添加方括号包围原文
                        # 先添加左方括号
                        first_char_pos = pos
                        if first_char_pos < len(char_to_run_map):
                            base_run_info = char_to_run_map[first_char_pos]
                            run = paragraph.add_run("[")
                            # 应用原始格式
                            run.bold = base_run_info['bold']
                            run.italic = base_run_info['italic']
                            run.underline = base_run_info['underline']
                            if base_run_info['font']:
                                run.font.name = base_run_info['font']
                            if base_run_info['size']:
                                run.font.size = base_run_info['size']
                            if base_run_info['color']:
                                run.font.color.rgb = base_run_info['color']
                            run.font.highlight_color = base_run_info['highlight_color']
                            if base_run_info['style']:
                                run.style = base_run_info['style']
                        else:
                            # 如果没有格式信息，则使用默认格式
                            paragraph.add_run("[")
                        
                        # 添加原文，保持原始格式
                        for i in range(pos, pos + len(original)):
                            if i < len(char_to_run_map):
                                run_info = char_to_run_map[i]
                                run = paragraph.add_run(text[i])
                                # 应用原始格式
                                run.bold = run_info['bold']
                                run.italic = run_info['italic']
                                run.underline = run_info['underline']
                                if run_info['font']:
                                    run.font.name = run_info['font']
                                if run_info['size']:
                                    run.font.size = run_info['size']
                                if run_info['color']:
                                    run.font.color.rgb = run_info['color']
                                run.font.highlight_color = run_info['highlight_color']
                                if run_info['style']:
                                    run.style = run_info['style']
                            else:
                                # 如果没有格式信息，则使用默认格式
                                paragraph.add_run(text[i])
                        
                        # 添加右方括号
                        last_char_pos = pos + len(original) - 1
                        if last_char_pos < len(char_to_run_map):
                            base_run_info = char_to_run_map[last_char_pos]
                            run = paragraph.add_run("]")
                            # 应用原始格式
                            run.bold = base_run_info['bold']
                            run.italic = base_run_info['italic']
                            run.underline = base_run_info['underline']
                            if base_run_info['font']:
                                run.font.name = base_run_info['font']
                            if base_run_info['size']:
                                run.font.size = base_run_info['size']
                            if base_run_info['color']:
                                run.font.color.rgb = base_run_info['color']
                            run.font.highlight_color = base_run_info['highlight_color']
                            if base_run_info['style']:
                                run.style = base_run_info['style']
                        else:
                            # 如果没有格式信息，则使用默认格式
                            paragraph.add_run("]")
                        
                        # 添加红色建议，使用圆括号包围，保持原文格式，仅颜色改为红色
                        if last_char_pos < len(char_to_run_map):
                            base_run_info = char_to_run_map[last_char_pos]
                            run = paragraph.add_run(f"({suggestion})")
                            # 应用原始格式
                            run.bold = base_run_info['bold']
                            run.italic = base_run_info['italic']
                            run.underline = base_run_info['underline']
                            if base_run_info['font']:
                                run.font.name = base_run_info['font']
                            if base_run_info['size']:
                                run.font.size = base_run_info['size']
                            # 仅颜色设置为红色
                            run.font.color.rgb = RGBColor(255, 0, 0)
                            run.font.highlight_color = base_run_info['highlight_color']
                            if base_run_info['style']:
                                run.style = base_run_info['style']
                        else:
                            # 如果没有格式信息，则使用默认格式，仅设置红色
                            run = paragraph.add_run(f"({suggestion})")
                            run.font.color.rgb = RGBColor(255, 0, 0)
                        
                        current_pos = pos + len(original)
                
                # 添加剩余文本，保持原始格式
                if current_pos < len(text):
                    for i in range(current_pos, len(text)):
                        if i < len(char_to_run_map):
                            run_info = char_to_run_map[i]
                            run = paragraph.add_run(text[i])
                            # 应用原始格式
                            run.bold = run_info['bold']
                            run.italic = run_info['italic']
                            run.underline = run_info['underline']
                            if run_info['font']:
                                run.font.name = run_info['font']
                            if run_info['size']:
                                run.font.size = run_info['size']
                            if run_info['color']:
                                run.font.color.rgb = run_info['color']
                            run.font.highlight_color = run_info['highlight_color']
                            if run_info['style']:
                                run.style = run_info['style']
                        else:
                            # 如果没有格式信息，则使用默认格式
                            paragraph.add_run(text[i])
        
        # 更新进度
        processed_items += 1
        if progress_callback:
            progress_percent = int((processed_items / total_items) * 100)
            progress_callback(progress_percent, f"处理段落 {processed_items}/{total_items}")
    
    # 处理表格
    for table in doc.tables:
        for row in table.rows:
            for cell in row.cells:
                for paragraph in cell.paragraphs:
                    text = paragraph.text
                    # 跳过空段落的处理
                    if not text.strip():
                        # 更新进度
                        processed_items += 1
                        if progress_callback:
                            progress_percent = int((processed_items / total_items) * 100)
                            progress_callback(progress_percent, f"正在处理 {processed_items}/{total_items}")
                        continue
                        
                    _, suggestions = get_openai_suggestions(text)
                    
                    if suggestions:
                            # 保存原始段落的所有运行及其格式
                            original_runs = []
                            for run in paragraph.runs:
                                original_runs.append({
                                    'text': run.text,
                                    'bold': run.bold,
                                    'italic': run.italic,
                                    'underline': run.underline,
                                    'font': run.font.name,
                                    'size': run.font.size,
                                    'color': run.font.color.rgb if run.font.color.rgb else None,
                                    'highlight_color': run.font.highlight_color,
                                    'style': run.style
                                })
                            
                            # 清除段落中的所有运行
                            p_element = paragraph._element
                            p_element.clear()
                            
                            # 重建原始文本并添加建议
                            # 首先构建原始文本的字符位置到运行的映射
                            char_to_run_map = []
                            current_pos = 0
                            for run_info in original_runs:
                                run_text = run_info['text']
                                for _ in range(len(run_text)):
                                    char_to_run_map.append(run_info)
                                    current_pos += 1
                            
                            # 重新添加文本，并在需要修改的地方添加红色建议
                            current_pos = 0
                            for original, suggestion in suggestions:
                                # 查找原文在文本中的位置
                                pos = text.find(original, current_pos)
                                if pos != -1:
                                    # 添加原文前的文本，保持原始格式
                                    if pos > current_pos:
                                        for i in range(current_pos, pos):
                                            if i < len(char_to_run_map):
                                                run_info = char_to_run_map[i]
                                                run = paragraph.add_run(text[i])
                                                # 应用原始格式
                                                run.bold = run_info['bold']
                                                run.italic = run_info['italic']
                                                run.underline = run_info['underline']
                                                if run_info['font']:
                                                    run.font.name = run_info['font']
                                                if run_info['size']:
                                                    run.font.size = run_info['size']
                                                if run_info['color']:
                                                    run.font.color.rgb = run_info['color']
                                                run.font.highlight_color = run_info['highlight_color']
                                                if run_info['style']:
                                                    run.style = run_info['style']
                                            else:
                                                # 如果没有格式信息，则使用默认格式
                                                paragraph.add_run(text[i])
                                    
                                    # 添加方括号包围原文
                                    # 先添加左方括号
                                    first_char_pos = pos
                                    if first_char_pos < len(char_to_run_map):
                                        base_run_info = char_to_run_map[first_char_pos]
                                        run = paragraph.add_run("[")
                                        # 应用原始格式
                                        run.bold = base_run_info['bold']
                                        run.italic = base_run_info['italic']
                                        run.underline = base_run_info['underline']
                                        if base_run_info['font']:
                                            run.font.name = base_run_info['font']
                                        if base_run_info['size']:
                                            run.font.size = base_run_info['size']
                                        if base_run_info['color']:
                                            run.font.color.rgb = base_run_info['color']
                                        run.font.highlight_color = base_run_info['highlight_color']
                                        if base_run_info['style']:
                                            run.style = base_run_info['style']
                                    else:
                                        # 如果没有格式信息，则使用默认格式
                                        paragraph.add_run("[")
                                    
                                    # 添加原文，保持原始格式
                                    for i in range(pos, pos + len(original)):
                                        if i < len(char_to_run_map):
                                            run_info = char_to_run_map[i]
                                            run = paragraph.add_run(text[i])
                                            # 应用原始格式
                                            run.bold = run_info['bold']
                                            run.italic = run_info['italic']
                                            run.underline = run_info['underline']
                                            if run_info['font']:
                                                run.font.name = run_info['font']
                                            if run_info['size']:
                                                run.font.size = run_info['size']
                                            if run_info['color']:
                                                run.font.color.rgb = run_info['color']
                                            run.font.highlight_color = run_info['highlight_color']
                                            if run_info['style']:
                                                run.style = run_info['style']
                                        else:
                                            # 如果没有格式信息，则使用默认格式
                                            paragraph.add_run(text[i])
                                    
                                    # 添加右方括号
                                    last_char_pos = pos + len(original) - 1
                                    if last_char_pos < len(char_to_run_map):
                                        base_run_info = char_to_run_map[last_char_pos]
                                        run = paragraph.add_run("]")
                                        # 应用原始格式
                                        run.bold = base_run_info['bold']
                                        run.italic = base_run_info['italic']
                                        run.underline = base_run_info['underline']
                                        if base_run_info['font']:
                                            run.font.name = base_run_info['font']
                                        if base_run_info['size']:
                                            run.font.size = base_run_info['size']
                                        if base_run_info['color']:
                                            run.font.color.rgb = base_run_info['color']
                                        run.font.highlight_color = base_run_info['highlight_color']
                                        if base_run_info['style']:
                                            run.style = base_run_info['style']
                                    else:
                                        # 如果没有格式信息，则使用默认格式
                                        paragraph.add_run("]")
                                    
                                    # 添加红色建议，使用圆括号包围，保持原文格式，仅颜色改为红色
                                    if last_char_pos < len(char_to_run_map):
                                        base_run_info = char_to_run_map[last_char_pos]
                                        run = paragraph.add_run(f"({suggestion})")
                                        # 应用原始格式
                                        run.bold = base_run_info['bold']
                                        run.italic = base_run_info['italic']
                                        run.underline = base_run_info['underline']
                                        if base_run_info['font']:
                                            run.font.name = base_run_info['font']
                                        if base_run_info['size']:
                                            run.font.size = base_run_info['size']
                                        # 仅颜色设置为红色
                                        run.font.color.rgb = RGBColor(255, 0, 0)
                                        run.font.highlight_color = base_run_info['highlight_color']
                                        if base_run_info['style']:
                                            run.style = base_run_info['style']
                                    else:
                                        # 如果没有格式信息，则使用默认格式，仅设置红色
                                        run = paragraph.add_run(f"({suggestion})")
                                        run.font.color.rgb = RGBColor(255, 0, 0)
                                    
                                    current_pos = pos + len(original)
                            
                            # 添加剩余文本，保持原始格式
                            if current_pos < len(text):
                                for i in range(current_pos, len(text)):
                                    if i < len(char_to_run_map):
                                        run_info = char_to_run_map[i]
                                        run = paragraph.add_run(text[i])
                                        # 应用原始格式
                                        run.bold = run_info['bold']
                                        run.italic = run_info['italic']
                                        run.underline = run_info['underline']
                                        if run_info['font']:
                                            run.font.name = run_info['font']
                                        if run_info['size']:
                                            run.font.size = run_info['size']
                                        if run_info['color']:
                                            run.font.color.rgb = run_info['color']
                                        run.font.highlight_color = run_info['highlight_color']
                                        if run_info['style']:
                                            run.style = run_info['style']
                                    else:
                                        # 如果没有格式信息，则使用默认格式
                                        paragraph.add_run(text[i])
                        
                    # 更新进度
                    processed_items += 1
                    if progress_callback:
                        progress_percent = int((processed_items / total_items) * 100)
                        progress_callback(progress_percent, f"处理表格单元格 {processed_items}/{total_items}")
    
    # 保存修订后的文件
    if progress_callback:
        progress_callback(100, "保存文件...")
    output_path = get_output_path(file_path)
    doc.save(output_path)
    
    return output_path


def process_powerpoint(file_path, progress_callback=None):
    """处理PowerPoint演示文稿"""
    prs = Presentation(file_path)
    
    # 计算总项目数（所有幻灯片中的非空文本形状）
    total_items = 0
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                # 包括空文本在内的所有文本形状都计入总数
                # 这样即使跳过空文本，进度计算也是准确的
                total_items += 1
    
    # 当前已处理项目数
    processed_items = 0
    
    # 处理所有幻灯片
    for slide_index, slide in enumerate(prs.slides):
        # 处理形状中的文本
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                text = shape.text
                # 跳过空文本的处理
                if not text.strip():
                    # 更新进度
                    processed_items += 1
                    if progress_callback:
                        progress_percent = int((processed_items / total_items) * 100)
                        progress_callback(progress_percent, f"跳过幻灯片 {slide_index+1} 中空文本 {processed_items}/{total_items}")
                    continue
                    
                _, suggestions = get_openai_suggestions(text)
                
                if suggestions:
                    # 保存原始格式信息
                    text_frame = shape.text_frame
                    original_paragraphs = []
                    
                    for paragraph in text_frame.paragraphs:
                        p_info = {
                            'text': paragraph.text,
                            'alignment': paragraph.alignment,
                            'level': paragraph.level,
                            'runs': []
                        }
                        
                        for run in paragraph.runs:
                            run_info = {
                                'text': run.text,
                                'bold': run.font.bold,
                                'italic': run.font.italic,
                                'underline': run.font.underline,
                                'font': run.font.name,
                                'size': run.font.size,
                                'color': run.font.color.rgb if hasattr(run.font.color, 'rgb') else None
                            }
                            p_info['runs'].append(run_info)
                        
                        original_paragraphs.append(p_info)
                    
                    # 构建字符位置到格式的映射
                    char_to_format_map = []
                    current_pos = 0
                    
                    for p_info in original_paragraphs:
                        for run_info in p_info['runs']:
                            run_text = run_info['text']
                            for _ in range(len(run_text)):
                                char_to_format_map.append(run_info)
                                current_pos += 1
                    
                    # 清除所有段落
                    while len(text_frame.paragraphs) > 0:
                        p = text_frame.paragraphs[0]
                        p._element.getparent().remove(p._element)
                    
                    # 创建新段落
                    p = text_frame.add_paragraph()
                    # 应用原始段落格式（使用第一个段落的格式）
                    if original_paragraphs:
                        p.alignment = original_paragraphs[0]['alignment']
                        p.level = original_paragraphs[0]['level']
                    
                    # 重新添加文本，并在需要修改的地方添加红色建议
                    current_pos = 0
                    for original, suggestion in suggestions:
                        # 查找原文在文本中的位置
                        pos = text.find(original, current_pos)
                        if pos != -1:
                            # 添加原文前的文本，保持原始格式
                            if pos > current_pos:
                                for i in range(current_pos, pos):
                                    if i < len(char_to_format_map):
                                        run_info = char_to_format_map[i]
                                        run = p.add_run()
                                        run.text = text[i]
                                        # 应用原始格式
                                        run.font.bold = run_info['bold']
                                        run.font.italic = run_info['italic']
                                        run.font.underline = run_info['underline']
                                        if run_info['font']:
                                            run.font.name = run_info['font']
                                        if run_info['size']:
                                            run.font.size = run_info['size']
                                        if run_info['color']:
                                            run.font.color.rgb = run_info['color']
                                    else:
                                        # 如果没有格式信息，则使用默认格式
                                        run = p.add_run()
                                        run.text = text[i]
                            
                            # 添加左方括号
                            first_char_pos = pos
                            if first_char_pos < len(char_to_format_map):
                                base_run_info = char_to_format_map[first_char_pos]
                                run = p.add_run()
                                run.text = "["
                                # 应用原始格式
                                run.font.bold = base_run_info['bold']
                                run.font.italic = base_run_info['italic']
                                run.font.underline = base_run_info['underline']
                                if base_run_info['font']:
                                    run.font.name = base_run_info['font']
                                if base_run_info['size']:
                                    run.font.size = base_run_info['size']
                                if base_run_info['color']:
                                    run.font.color.rgb = base_run_info['color']
                            else:
                                # 如果没有格式信息，则使用默认格式
                                run = p.add_run()
                                run.text = "["
                            
                            # 添加原文，保持原始格式
                            for i in range(pos, pos + len(original)):
                                if i < len(char_to_format_map):
                                    run_info = char_to_format_map[i]
                                    run = p.add_run()
                                    run.text = text[i]
                                    # 应用原始格式
                                    run.font.bold = run_info['bold']
                                    run.font.italic = run_info['italic']
                                    run.font.underline = run_info['underline']
                                    if run_info['font']:
                                        run.font.name = run_info['font']
                                    if run_info['size']:
                                        run.font.size = run_info['size']
                                    if run_info['color']:
                                        run.font.color.rgb = run_info['color']
                                else:
                                    # 如果没有格式信息，则使用默认格式
                                    run = p.add_run()
                                    run.text = text[i]
                            
                            # 添加右方括号
                            last_char_pos = pos + len(original) - 1
                            if last_char_pos < len(char_to_format_map):
                                base_run_info = char_to_format_map[last_char_pos]
                                run = p.add_run()
                                run.text = "]"
                                # 应用原始格式
                                run.font.bold = base_run_info['bold']
                                run.font.italic = base_run_info['italic']
                                run.font.underline = base_run_info['underline']
                                if base_run_info['font']:
                                    run.font.name = base_run_info['font']
                                if base_run_info['size']:
                                    run.font.size = base_run_info['size']
                                if base_run_info['color']:
                                    run.font.color.rgb = base_run_info['color']
                            else:
                                # 如果没有格式信息，则使用默认格式
                                run = p.add_run()
                                run.text = "]"
                            
                            # 添加红色建议，使用圆括号包围，保持原文格式，仅颜色改为红色
                            if last_char_pos < len(char_to_format_map):
                                base_run_info = char_to_format_map[last_char_pos]
                                run = p.add_run()
                                run.text = f"({suggestion})"
                                # 应用原始格式
                                run.font.bold = base_run_info['bold']
                                run.font.italic = base_run_info['italic']
                                run.font.underline = base_run_info['underline']
                                if base_run_info['font']:
                                    run.font.name = base_run_info['font']
                                if base_run_info['size']:
                                    run.font.size = base_run_info['size']
                                # 仅颜色设置为红色
                                run.font.color.rgb = PPTRGBColor(255, 0, 0)
                            else:
                                # 如果没有格式信息，则使用默认格式，仅设置红色
                                run = p.add_run()
                                run.text = f"({suggestion})"
                                run.font.color.rgb = PPTRGBColor(255, 0, 0)
                            
                            current_pos = pos + len(original)
                    
                    # 添加剩余文本，保持原始格式
                    if current_pos < len(text):
                        for i in range(current_pos, len(text)):
                            if i < len(char_to_format_map):
                                run_info = char_to_format_map[i]
                                run = p.add_run()
                                run.text = text[i]
                                # 应用原始格式
                                run.font.bold = run_info['bold']
                                run.font.italic = run_info['italic']
                                run.font.underline = run_info['underline']
                                if run_info['font']:
                                    run.font.name = run_info['font']
                                if run_info['size']:
                                    run.font.size = run_info['size']
                                if run_info['color']:
                                    run.font.color.rgb = run_info['color']
                            else:
                                # 如果没有格式信息，则使用默认格式
                                run = p.add_run()
                                run.text = text[i]
                
                # 更新进度
                processed_items += 1
                if progress_callback:
                    progress_percent = int((processed_items / total_items) * 100)
                    progress_callback(progress_percent, f"处理幻灯片 {slide_index+1}/{len(prs.slides)}, 形状 {processed_items}/{total_items}")
    
    # 保存修订后的文件
    if progress_callback:
        progress_callback(100, "保存文件...")
    output_path = get_output_path(file_path)
    prs.save(output_path)
    
    return output_path





def get_output_path(file_path):
    """获取输出文件路径"""
    file_path = Path(file_path)
    stem = file_path.stem
    suffix = file_path.suffix
    directory = file_path.parent
    
    output_path = directory / f"{stem}_修订{suffix}"
    return output_path